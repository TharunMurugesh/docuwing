import re
from pathlib import Path
from sqlalchemy import delete, select
from sqlalchemy.ext.asyncio import AsyncSession
from app.models import Document, DocumentChunk


class KnowledgeEngine:
    """Deterministic parsing, structure-aware chunking, and project-scoped hybrid retrieval."""
    def __init__(self, session: AsyncSession): self.session = session
    async def ingest(self, document: Document, progress) -> int:
        await progress(0.1, "parsing")
        text, pages = self._extract(Path(document.storage_path), document.mime_type)
        await progress(0.4, "chunking")
        chunks = self._chunk(text)
        await self.session.execute(delete(DocumentChunk).where(DocumentChunk.document_id == document.id))
        for sequence, (content, metadata) in enumerate(chunks):
            self.session.add(DocumentChunk(document_id=document.id, sequence=sequence, content=content, metadata_json=metadata))
        document.page_count = pages; document.status = "ready"
        await self.session.commit(); await progress(1, "ready")
        return len(chunks)
    def _extract(self, path: Path, mime: str) -> tuple[str, int]:
        suffix = path.suffix.lower()
        if suffix in {".txt", ".md", ".csv"}: return path.read_text(encoding="utf-8", errors="replace"), 1
        if suffix == ".pdf":
            import fitz
            pdf = fitz.open(path); text = "\n\n".join(page.get_text() for page in pdf)
            if text.strip(): return text, len(pdf)
            return self._ocr_pdf(pdf), len(pdf)
        if suffix == ".docx":
            from docx import Document as Docx
            doc = Docx(path); return "\n".join(p.text for p in doc.paragraphs), 1
        if suffix == ".pptx":
            from pptx import Presentation
            deck = Presentation(path); return "\n\n".join("\n".join(s.text for s in slide.shapes if hasattr(s, "text")) for slide in deck.slides), len(deck.slides)
        if suffix == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(path, read_only=True, data_only=True); return "\n\n".join("# " + ws.title + "\n" + "\n".join(" | ".join(str(v or "") for v in row) for row in ws.iter_rows(values_only=True)) for ws in wb.worksheets), len(wb.worksheets)
        if suffix in {".png", ".jpg", ".jpeg"}:
            from PIL import Image
            import pytesseract
            return pytesseract.image_to_string(Image.open(path)), 1
        raise ValueError(f"Unsupported document type: {suffix}")

    def _ocr_pdf(self, pdf) -> str:
        import io
        from PIL import Image
        import pytesseract
        pages = []
        for page in pdf:
            pixmap = page.get_pixmap(dpi=300)
            pages.append(pytesseract.image_to_string(Image.open(io.BytesIO(pixmap.tobytes("png")))))
        return "\n\n".join(pages)
    def _chunk(self, text: str, size: int = 1800) -> list[tuple[str, dict]]:
        blocks = [b.strip() for b in re.split(r"\n\s*\n", text) if b.strip()]
        result: list[tuple[str, dict]] = []; current = ""; heading = ""
        for block in blocks:
            if block.startswith("#"): heading = block.split("\n", 1)[0].lstrip("# ")
            if current and len(current) + len(block) > size:
                result.append((current, {"heading": heading})); current = ""
            current += ("\n\n" if current else "") + block
        if current: result.append((current, {"heading": heading}))
        return result
    async def retrieve(self, project_id: str, query: str, limit: int = 8) -> list[dict]:
        words = {w.lower() for w in re.findall(r"\w+", query) if len(w) > 2}
        rows = (await self.session.execute(select(DocumentChunk, Document).join(Document).where(Document.project_id == project_id))).all()
        ranked = []
        for chunk, document in rows:
            terms = set(re.findall(r"\w+", chunk.content.lower()))
            score = len(words & terms) / max(len(words), 1)
            if score: ranked.append((score, chunk, document))
        ranked.sort(key=lambda item: item[0], reverse=True)
        return [{"chunk_id": c.id, "document_id": d.id, "document": d.filename, "sequence": c.sequence, "content": c.content, "metadata": c.metadata_json, "score": score} for score, c, d in ranked[:limit]]
